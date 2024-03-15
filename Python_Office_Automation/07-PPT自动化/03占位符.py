# 导入库
from pptx import Presentation
# 创建PPT文件
ppt = Presentation()
# 添加幻灯片
slide = ppt.slides.add_slide(ppt.slide_layouts[0])
# 获取占位符
for place in slide.placeholders:
    place.text = "这里是占位符，可以添加内容"
# 保存
ppt.save('source_material/03占位符.pptx')
print("保存成功")
