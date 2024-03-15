from pptx import Presentation
ppt = Presentation('source_material/13删除.pptx')
page = list(ppt.slides._sldIdLst) # 获取页的列表
print(len(page))
ppt.slides._sldIdLst.remove(page[0]) # 指定删除页
# ppt.slides._sldIdLst.remove(page[len(page)-1]) # 删除最后一页
ppt.save('source_material/13删除【结果】.pptx')
print('删除成功')
