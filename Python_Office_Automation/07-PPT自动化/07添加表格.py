from pptx import Presentation
from pptx.util import Inches

ppt = Presentation()
slide = ppt.slides.add_slide(ppt.slide_layouts[5])
shapes = slide.shapes

shapes.title.text = '添加表格'
# 设置行数列数,与边界的距离以及表格的宽高
rows = cols = 4
top = Inches(2)
left = Inches(2.5)
width = Inches(5)
height = Inches(0.8)

table = shapes.add_table(rows, cols, left, top, width, height).table

#设置列的宽度
#table.columns[1].width = Inches(3.0)
#table.columns[2].width = Inches(3.0)

data = [
    ["姓名", "学校", "专业", "成绩"],
    ["张三", "郑州大学", "会计", 99],
    ["李四", "清华大学", "软件工程", 89],
    ["王五", "清华大学", "金融", 95],
]
 
for row in range(rows):
    for col in range(cols):
        table.cell(row, col).text = str(data[row][col])

ppt.save('source_material/07添加表格.pptx')
print("添加成功")
