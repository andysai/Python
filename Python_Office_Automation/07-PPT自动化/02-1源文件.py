from pptx import Presentation
# 加载源文件
ppt = Presentation("source_material/02-1源文件.pptx")
for slide in ppt.slides:
    print(slide)
    for shape in slide.shapes:
        print(shape.text)
        if shape.has_text_frame:
            p = shape.text_frame
            for pra in p.paragraphs:
                print(pra.text)

ppt.save("source_material/02-1源文件.pptx")
