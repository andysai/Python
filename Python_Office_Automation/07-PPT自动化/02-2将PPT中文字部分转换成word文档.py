from pptx import Presentation
from docx import Document
doc = Document()
ppt = Presentation("source_material/02-1源文件.pptx")
for slide in ppt.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:   # 判断是否存在文字
            text_box = shape.text_frame
            for para in text_box.paragraphs:
                if para.text != "":  # 只导出不为空的段落
                    doc.add_paragraph(para.text)
doc.save("source_material/02-1源文件.docx")
print("保存成功")               
