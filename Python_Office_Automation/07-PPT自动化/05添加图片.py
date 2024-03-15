# 导入库
from pptx import Presentation
from pptx.util import Inches
#图片路径及名称
img_path = 'source_material/pic.jpg'
# 创建ppt文件
ppt = Presentation()
# 创建幻灯片
slide = ppt.slides.add_slide(ppt.slide_layouts[6])
# 设置位置
left = top = Inches(0.5)
height = Inches(5.5)
picture = slide.shapes.add_picture(img_path, left, top, height=height)

ppt.save('source_material/05添加图片.pptx')
print("添加成功")
