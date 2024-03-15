from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION

# 创建幻灯片
ppt = Presentation()    # 初始化 ppt 文档
slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # slide幻灯片
chart = slide.shapes
title = slide.shapes.title
title.text = "GDP占比"

# 定义图表数据
x = ['郑州', '武汉', '北京', '上海']
y = [7848, 5273, 9089, 8659]
z = [3568, 2572, 3748, 4105]

chart_data = ChartData()
chart_data.categories = x    # 设置x轴
chart_data.add_series(name='上半年', values=z)
chart_data.add_series(name='下半年', values=y)

# 添加图表
left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(6)
# 簇状柱形图:chart_type=XL_CHART_TYPE.PIE 图表区的位置:x,y 图表的宽和高:cx,cy
creat_chart = chart.add_chart(chart_type=XL_CHART_TYPE.PIE, x=left, y=top, cx=width, cy=height, chart_data=chart_data)

set_chart = creat_chart.chart
layout_chart = set_chart.plots[0]
# 设置数据标签
layout_chart.has_data_labels = True    # 显示数据标签
data_head = layout_chart.data_labels    # 获取数据标签控制类
data_head.show_category_name = True    # 是否显示类别名称
data_head.show_value = False    # 是否显示值
data_head.show_percentage = True    # 是否显示百分比
data_head.number_format = '0.0%'    # 标签的数字格式
data_head.position = XL_LABEL_POSITION.INSIDE_END    # 标签位置
data_head.font.name = '宋体'
data_head.font.size = Pt(20)

# 保存
ppt.save('source_material/12饼图.pptx')
print("生成成功")
