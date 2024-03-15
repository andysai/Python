from pptx import Presentation
from pptx.util import Pt
ppt = Presentation()
slide = ppt.slides.add_slide(ppt.slide_layouts[1])
para = slide.shapes.placeholders
title = slide.shapes.title
title.text = "将进酒"
para_text = para[1].text_frame.add_paragraph()  # 在第二个shape中的文本框架中添加段落
para_text.text = '君不见行河之水天上来，奔流到海不复回。'  # 段落中文字内容
para_text.font.bold = True  # 文字加粗
para_text.font.italic = True  # 文字斜体
para_text.font.size = Pt(25)  # 文字大小
para_text.font.underline = False  # 文字无下划线
para_text.level = 1  # 段落的级别
ppt.save("source_material/06添加段落.pptx")
print("添加成功")
