from pptx import Presentation
from pptx.util import Inches
ppt = Presentation('source_material/工作证.pptx')
slide = ppt.slides.add_slide(ppt.slide_layouts[0])
# 获取占位符下标
#for place in slide.placeholders: # 遍历母版中的占位符
#    message = place.placeholder_format
#    place.text = f'{message.idx}'
#    print(place.text)
# 通过获取的下标填充数据
name = slide.placeholders[11]
gender = slide.placeholders[12]
ID = slide.placeholders[13]
top = Inches(2.5)
left = Inches(2.8)
height = Inches(1.5)
picture = slide.shapes.add_picture(f'source_material/张三.jpg', left, top, height=height)
name.text = '张三'
gender.text = '男'
ID.text = '202001'
ppt.save('source_material/09模板渲染—工作证.pptx')
print("保存成功")
