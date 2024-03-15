from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
ppt = Presentation()
slide = ppt.slides.add_slide(ppt.slide_layouts[5])
title = slide.shapes.title
title.text = "添加形状"
# 预设位置及大小
left, top, width, height = Inches(1), Inches(3), Inches(1.8), Inches(1)
# 在指定位置按预设值添加类型为CAN的形状
shape = slide.shapes.add_shape(MSO_SHAPE.CAN, left, top, width, height)
shape.text = 'Step 1'
for n in range(2, 6):
  left = left + width - Inches(0.3)
  shape = slide.shapes.add_shape(MSO_SHAPE.CAN, left, top, width, height)
  shape.text = 'Step{}'.format(n)
ppt.save('source_material/08添加图形.pptx')
print("添加成功")
