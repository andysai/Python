# 导入库
from pptx import Presentation
from pptx.util import Inches
# 创建PPT文件
ppt = Presentation()
# 创建幻灯片
slide = ppt.slides.add_slide(ppt.slide_layouts[0])
# 获取主题占位符
title = slide.shapes.title
title.text = "主题占位符"
# 获取副标题占位符
place = slide.placeholders[1]
place.text = "副标题占位符"
# 设置添加文本框的位置和宽高
top = left = Inches(0)
width = Inches(3)
height = Inches(2)
# 添加文本框
textbox = slide.shapes.add_textbox(top=top, left=left, width=width, height=height)
textbox.text = "添加的文本框"
# 保存
ppt.save("source_material/04添加文本框.pptx")
print("保存成功")











