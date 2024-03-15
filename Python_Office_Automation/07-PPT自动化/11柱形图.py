from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_DATA_LABEL_POSITION
# 创建幻灯片
ppt = Presentation()
slide = ppt.slides.add_slide(ppt.slide_layouts[6]) 
chart = slide.shapes

# 定义图表数据
x = ['郑州', '武汉', '北京', '上海']
y = [7848, 5273, 9089, 8659]
z = [3568, 2572, 3748, 4105]

chart_data = ChartData()
chart_data.categories = x    # 设置x轴
chart_data.add_series(name='上半年', values=z)
chart_data.add_series(name='下半年', values=y)

# 添加图表
left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(6)
# 簇状柱形图:chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED 图表区的位置:x,y 图表的宽和高:cx,cy
creat_chart = chart.add_chart(chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED, x=left, y=top, cx=width, cy=height, chart_data=chart_data)
# 设置标题，字体，颜色，位置等
get_chart = creat_chart.chart  # 从生成的图表中取出图表类
get_chart.chart_style = 10  # 图表整体颜色风格
get_chart.has_title = True  # 图表是否含有标题，默认为False
# get_chart.chart_title.text_frame.clear()  # 清除原标题
# new_paragraph = get_chart.chart_title.text_frame.add_paragraph()  # 添加一行新标题
# new_paragraph.text = 'GDP   单位：亿元'  # 新标题
# new_paragraph.font.size = Pt(15)  # 新标题字体大小

plot = get_chart.plots[0]  # 取图表中第一个plot
plot.has_data_labels = True  # 是否显示数据标签
data_labels = plot.data_labels  # 数据标签控制类
data_labels.font.size = Pt(13)  # 字体大小
data_labels.font.color.rgb = RGBColor(0, 0, 255)  # 字体颜色
data_labels.position = XL_DATA_LABEL_POSITION.INSIDE_END  # 字体位置

ppt.save('source_material/11柱形图.pptx')
print("生成成功")
