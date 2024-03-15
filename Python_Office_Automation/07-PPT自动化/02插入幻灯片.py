# 导入库
from pptx import Presentation
# 创建ppt文件
ppt = Presentation()
# 添加幻灯片，选择第二的母版
slide = ppt.slides.add_slide(ppt.slide_layouts[1])
# 保存
ppt.save('source_material/02添加幻灯片.pptx')
print("success")
