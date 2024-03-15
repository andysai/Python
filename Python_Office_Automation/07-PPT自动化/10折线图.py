from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION

# 创建幻灯片
ppt = Presentation()    # 初始化 ppt 文档
slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # slide幻灯片
chart = slide.shapes
title = slide.shapes.title
title.text = "GDP   单位：亿元"

# 定义图表数据
x = ['郑州', '武汉', '北京', '上海']
y = [7848, 5273, 9089, 8659]
z = [3568, 2572, 3748, 4105]

chart_data = ChartData()
chart_data.categories = x    # 设置x轴
chart_data.add_series(name='上半年', values=z)
chart_data.add_series(name='下半年', values=y)

# 添加图表
left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(6)
# 图表类型:chart_type=XL_CHART_TYPE.LINE 图表区的位置:x,y 图表的宽和高:cx,cy
creat_chart = chart.add_chart(chart_type=XL_CHART_TYPE.LINE, x=left, y=top, cx=width, cy=height, chart_data=chart_data)

ppt.save('source_material/10折线图.pptx')
print("生成成功")
